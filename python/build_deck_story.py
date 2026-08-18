"""Generates the article-companion .pptx for the Graph RAG / ArangoDB project --
same story, same shape (12 slides) as build_deck.py, but for a general reader:
minimal jargon, narrative framing, "why this is interesting" over implementation
detail. Content is drawn from the same source docs as the technical deck
(_project_step_by_step_plan.md, _step3_extraction_reference.md), just retold.

Re-run any time to regenerate from scratch. Layout helpers live in
deck_theme.py, shared with build_deck.py so both decks look like a matched set.
"""

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from deck_theme import (
    NAVY, TEAL, LIGHT_TEAL, DARK_TEXT, WHITE, GRAY, LIGHT_BLUE_TEXT,
    new_presentation, add_slide, set_background, add_textbox, add_bullets,
    add_header, add_footer, add_box, add_plain_box, add_arrow, add_callout,
)

FOOTER_LABEL = "From 187 Filings to a Knowledge Graph"

prs, BLANK = new_presentation()

# ============================================================================
# Slide 1 -- Title
# ============================================================================
s = add_slide(prs, BLANK)
set_background(s, NAVY)
add_textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.2),
            "Teaching a Computer to Read a Case File", size=40, color=WHITE, bold=True)
add_textbox(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.8),
            "How 187 real court filings became a map you can ask questions of",
            size=20, color=TEAL, bold=False)
add_textbox(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.6),
            "A case study in building with AI, one careful step at a time", size=15, color=LIGHT_BLUE_TEXT, italic=True)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.35), Inches(2.2), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = TEAL
line.line.fill.background()
line.shadow.inherit = False

# ============================================================================
# Slide 2 -- The Problem
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "A Mountain of Paperwork", kicker="Why bother?")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "187 real court documents from a single federal case — close to 200 megabytes of dense legal filings",
    "Somewhere in there: names, roles, dates, and relationships — but buried in paragraphs, not organized anywhere",
    "A person could read it all eventually — but answering \"everyone this witness is connected to\" would mean re-reading everything, every single time",
    "The goal: turn a pile of documents into something you can actually ask questions of, with every answer traceable back to the page it came from",
    "One rule from day one: nothing in these filings gets handed to an outside AI service without a deliberate decision — the files likely name real people who aren't public figures",
], size=18, space_after=18)
add_footer(s, 2, FOOTER_LABEL)

# ============================================================================
# Slide 3 -- The Big Idea
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "A Map, and a Sense of Meaning", kicker="Two kinds of memory")

col_w = Inches(5.6)
col_h = Inches(3.0)
top = Inches(1.7)
add_plain_box(s, Inches(0.7), top, col_w, col_h, NAVY)
add_plain_box(s, Inches(6.9), top, col_w, col_h, TEAL)

add_textbox(s, Inches(1.0), Inches(1.85), col_w - Inches(0.6), Inches(0.4),
            "The Map", size=20, color=WHITE, bold=True)
add_textbox(s, Inches(1.0), Inches(2.25), col_w - Inches(0.6), Inches(0.3),
            "who's connected to whom", size=13, color=LIGHT_BLUE_TEXT, italic=True)
add_bullets(s, Inches(1.0), Inches(2.65), col_w - Inches(0.6), Inches(2), [
    "Keeps track of people, organizations, and courts as real, connected records",
    "Answers \"who's connected to whom, and how\" directly, instead of re-reading everything",
], size=14, color=WHITE, space_after=10)

add_textbox(s, Inches(7.2), Inches(1.85), col_w - Inches(0.6), Inches(0.4),
            "The Sense of Meaning", size=20, color=WHITE, bold=True)
add_textbox(s, Inches(7.2), Inches(2.25), col_w - Inches(0.6), Inches(0.3),
            "finds ideas, not just exact words", size=13, color=LIGHT_TEAL, italic=True)
add_bullets(s, Inches(7.2), Inches(2.65), col_w - Inches(0.6), Inches(2), [
    "Can find a passage that means the same thing even if it uses completely different words",
    "This is what lets you ask a plain-English question later and get relevant passages back",
], size=14, color=WHITE, space_after=10)

add_textbox(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(0.4),
            "How it's built", size=16, color=NAVY, bold=True)
add_bullets(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.8), [
    "A local AI model does the reading — nothing about these documents leaves the machine",
    "One piece of software handles reading the PDFs; a separate piece runs the web interface and shows live progress while files upload",
], size=15)
add_footer(s, 3, FOOTER_LABEL)

# ============================================================================
# Slide 4 -- The Journey at a Glance (diagram)
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "The Journey, at a Glance", kicker="Six stages, repeated for every document")

steps = [
    ("Read & Slice", "PDF → bite-sized pieces"),
    ("Spot the Players", "who's mentioned, and how"),
    ("Find Matches", "same person again?"),
    ("Double-Check", "AI reviews the maybes"),
    ("Update the Map", "write it down for good"),
    ("Ask Questions", "get sourced answers"),
]
box_w = Inches(1.8)
box_h = Inches(1.15)
arrow_w = Inches(0.3)
start_x = Inches(0.5)
y = Inches(2.4)
x = start_x
for i, (label, sub) in enumerate(steps):
    fill = NAVY if i % 2 == 0 else TEAL
    add_box(s, x, y, box_w, box_h, label, sub, fill=fill, size=13)
    x = Emu(x + box_w)
    if i < len(steps) - 1:
        add_arrow(s, x, Emu(y + box_h // 2 - Emu(Inches(0.12))), arrow_w, Inches(0.24))
        x = Emu(x + arrow_w)

add_bullets(s, Inches(0.7), Inches(4.1), Inches(11.9), Inches(2.6), [
    "\"Read & Slice\" and \"Spot the Players\" are built and running against the real documents today",
    "\"Find Matches,\" \"Double-Check,\" and \"Update the Map\" are fully planned, not yet built",
    "\"Ask Questions\" comes last, on purpose — the map has to exist before anyone can query it",
    "At no point does the system try to hold all 187 documents in the AI's head at once — each stage works on one small piece, and a database remembers everything permanently, the way a researcher keeps a notebook instead of memorizing a library",
], size=15, space_after=10)
add_footer(s, 4, FOOTER_LABEL)

# ============================================================================
# Slide 5 -- Read & Slice
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Turning a Wall of Text into Bite-Sized Pieces", kicker="Read & Slice")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "AI models can only \"hold\" so much text in their head at once — feed them too much and quality drops fast",
    "So every document gets sliced into overlapping pieces, a few hundred words each, like index cards with a little overlap so nothing gets cut off mid-thought",
    "Each slice gets a \"fingerprint\" that captures what it means, not just the exact words used — this is what powers the later \"sense of meaning\" search",
    "Every slice is saved in two places and cross-checked against each other, so nothing quietly goes missing",
    "Re-uploading the same file safely replaces its old pieces instead of creating duplicates",
    ("Done and verified: 187 documents → 3,575 slices, confirmed complete", 0),
], size=17, space_after=14)
add_footer(s, 5, FOOTER_LABEL)

# ============================================================================
# Slide 6 -- Spot the Players
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Teaching the AI to Spot Who's Who", kicker="Spot the Players")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "Some facts don't need a guess at all: dates and case numbers follow a strict, predictable format, so simple pattern-matching finds them with zero room for error — no AI involved",
    "For the harder part — spotting people, organizations, and courts — a small AI model reads each slice and is given exactly one narrow job",
    "It's told to note someone's role (like \"witness\" or \"defendant\") only when the text says so directly, right next to their name — never to guess or assume",
    "Anything that's just a citation to another case or document gets automatically thrown out, checked two different ways",
    "An early version tried to have the AI do everything at once — people, roles, relationships, dates — and it started guessing instead of admitting uncertainty. The fix was to narrow its job, not to write a cleverer instruction",
], size=17, space_after=14)
add_footer(s, 6, FOOTER_LABEL)

# ============================================================================
# Slide 7 -- The time the AI invented a judge
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "The Time the AI Invented a Judge", kicker="A near miss")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(3.2), [
    "While testing, the AI read a case number ending in the letters \"LAP\" — a judge's initials, tacked on as a label — and confidently reported a person named \"Judge LAP\"",
    "It wasn't a typo or a fluke. It was the AI doing exactly what AI models do: finding a pattern and running with it",
    "The fix wasn't a cleverer instruction — instructions alone don't reliably stop a model that's allowed to interpret. The fix was to take that decision away from the AI entirely and hand it to simple, literal pattern-matching instead",
], size=17, space_after=16)
add_callout(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.6),
            "If a piece of information has one right answer and a fixed shape, don't ask an AI to guess it — write a rule instead.")
add_footer(s, 7, FOOTER_LABEL)

# ============================================================================
# Slide 8 -- Learning to live with "probably"
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Learning to Live with “Probably”", kicker="No AI is perfectly consistent")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "Ask the same AI the same question twice, and the answer can come back slightly different — even with settings meant to force consistency",
    "Rather than chase perfect consistency, this project treats \"not fully sure\" as an expected outcome, not a bug to eliminate",
    "When the AI genuinely can't tell if two mentions are the same person, it doesn't guess — it flags the pair for a human to look at later",
    "Wrongly merging two real people into one record is treated as worse than leaving them separate for now and sorting it out later",
    "Building a map like this by hand, one document at a time, would take an enormous amount of time and still wouldn't be flawless — so working in percentages, with people reviewing the genuinely uncertain cases, is the deliberate choice here",
], size=17, space_after=14)
add_footer(s, 8, FOOTER_LABEL)

# ============================================================================
# Slide 9 -- Find Matches / Double-Check / Update the Map
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Is “J. Smith” the Same as “John Smith”?", kicker="One person, many names")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "The same person can be mentioned dozens of different ways across 187 documents: full name, initials, just their role (\"defense counsel\"), a nickname",
    "Two checks run side by side: one compares spelling directly, the other compares meaning — so \"defense counsel\" and \"Ms. Walz\" can still be linked if they show up in similar surroundings",
    "Anything that looks like a plausible match gets a second opinion from the AI, which weighs the surrounding facts and decides: same person, different person, or genuinely unclear",
    "Confirmed matches update one shared record instead of piling up duplicates; a genuinely new person gets their own new entry",
    "This stage is fully designed, but not yet built — it's next in line once the current reading pass finishes",
], size=17, space_after=14)
add_footer(s, 9, FOOTER_LABEL)

# ============================================================================
# Slide 10 -- Ask Questions
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Someday: Just Ask", kicker="The payoff, later")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "Once the map and the meaning-index are built, the plan is to let someone type a plain-English question and get a real answer",
    "The system finds the most relevant passages and the people/entities connected to them, and has the AI compose an answer grounded in what it actually found",
    "Every answer traces back to the real filing it came from — nothing invented, nothing unsourced",
    "This stage comes last, on purpose: the whole point of everything before it is building a trustworthy map first, so the questions asked of it have real answers underneath them",
], size=19, space_after=18)
add_footer(s, 10, FOOTER_LABEL)

# ============================================================================
# Slide 11 -- Where Things Stand
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Where Things Stand", kicker="As of this deck")
stats = [
    ("187", "documents read in"),
    ("3,575", "slices created,\ndouble-checked complete"),
    ("714 / 3,575", "slices scanned\nfor people so far"),
]
box_w = Inches(3.7)
gap = Inches(0.4)
x = Inches(0.7)
for val, lbl in stats:
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), box_w, Inches(1.9))
    shp.fill.solid()
    shp.fill.fore_color.rgb = NAVY
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = val
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = TEAL
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = lbl
    r2.font.size = Pt(13)
    r2.font.color.rgb = WHITE
    x = Emu(x + box_w + gap)

add_bullets(s, Inches(0.7), Inches(4.1), Inches(11.9), Inches(2.6), [
    "Reading and slicing every document: done and verified",
    "Spotting the people, organizations, and courts: built and proven; currently running on the full set of documents",
    "Matching, double-checking, and updating the map: fully planned, no code written yet",
    "Letting someone ask questions of it: a later phase, not designed in detail yet",
], size=16, space_after=10)
add_footer(s, 11, FOOTER_LABEL)

# ============================================================================
# Slide 12 -- What's Next
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "What's Next", kicker="Closing")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "Finish scanning every document for people, organizations, and courts",
    "Build the matching stage — catching the same person mentioned under different names",
    "Build the AI double-check stage for anything the matching stage is unsure about",
    "Build the stage that actually writes confirmed matches into the map",
    "Later: a simple review page so a human can resolve the flagged “maybe the same person” cases on their own time",
    "Later still: teach the system to spot relationships between people, not just the people themselves, and let someone finally ask it questions directly",
], size=18, space_after=16)
add_footer(s, 12, FOOTER_LABEL)

# ============================================================================
out_path = "/Users/budubasa/Documents/N_Claude_Projects/arango_graph_rag_llm/Graph_RAG_LLM_Story.pptx"
prs.save(out_path)
print(f"Saved {out_path} ({len(prs.slides)} slides)")
