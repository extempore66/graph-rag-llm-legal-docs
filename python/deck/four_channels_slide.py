#!/usr/bin/env python3
"""
The four-channel fusion diagram, as a single editable PowerPoint slide.

Native shapes and connectors -- not an exported image -- so every box can be
moved, recoloured or deleted in PowerPoint. Mirrors the SVG version, with one
deliberate difference: the layout is mapped into the slide with a WIDER
horizontal scale than vertical. A top-to-bottom flow on a 16:9 slide is always
height-bound, and stretching the horizontal axis fills the space that would
otherwise be empty margin, while giving the box labels more room per line.

Caption and footnote go to the speaker notes rather than onto the face -- the
slide has to read from six metres away.

Usage:
    python/.venv/bin/python python/deck/four_channels_slide.py <out.pptx>
"""

import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# palette -- same values as the deck and the SVG
CARD = RGBColor(0xFB, 0xFA, 0xF8)
CARD2 = RGBColor(0xF4, 0xF2, 0xEE)
INK = RGBColor(0x14, 0x17, 0x1B)
INK_SOFT = RGBColor(0x4A, 0x51, 0x58)
INK_FAINT = RGBColor(0x83, 0x8B, 0x93)
RULE = RGBColor(0xD8, 0xD5, 0xCD)
SEAL = RGBColor(0x1E, 0x3A, 0x5F)
SEAL_WASH = RGBColor(0xE7, 0xEB, 0xF1)
OXIDE = RGBColor(0xA6, 0x38, 0x2C)
OXIDE_WASH = RGBColor(0xF6, 0xE9, 0xE7)

SERIF, SANS, MONO = "Georgia", "Helvetica Neue", "Menlo"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

# The SVG's coordinate space, mapped onto the slide's drawing area.
SX = 11.6 / 1000.0          # inches per svg x-unit
SY = 5.62 / 690.0           # inches per svg y-unit
OX, OY = 0.87, 1.34         # top-left of the drawing area, inches


def X(u):
    return Inches(OX + u * SX)


def Y(u):
    return Inches(OY + u * SY)


def W(u):
    return Inches(u * SX)


def H(u):
    return Inches(u * SY)


def arrowhead(conn):
    """
    python-pptx has no API for line ends, so the triangle is added directly to
    the drawing XML. Appended last because a:ln's children are ordered and
    tailEnd sits after the fill and dash elements we set above it.
    """
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"})
    ln.append(tail)


def line(slide, x1, y1, x2, y2, *, head=True, color=INK_FAINT):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, X(x1), Y(y1), X(x2), Y(y2))
    c.line.color.rgb = color
    c.line.width = Pt(1)
    if head:
        arrowhead(c)
    return c


def box(slide, x, y, w, h, *, fill, stroke, radius=0.05):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, X(x), Y(y), W(w), H(h))
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = stroke
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def lines_in(shape, rows):
    """rows: (text, font, size, color, bold) -- one centred paragraph each."""
    tf = shape.text_frame
    for i, (text, font, size, color, bold) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(1)
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold


def label(slide, cx, y, text, *, font=SANS, size=9, color=INK_SOFT, bold=False,
          width=300):
    tb = slide.shapes.add_textbox(X(cx - width / 2), Y(y), W(width), Inches(0.22))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def build(path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = CARD

    # ---- header -----------------------------------------------------------
    tb = slide.shapes.add_textbox(Inches(0.62), Inches(0.36), Inches(9), Inches(0.24))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "THE SOLUTION · RETRIEVAL"
    r.font.name, r.font.size, r.font.color.rgb = MONO, Pt(9), SEAL

    tb = slide.shapes.add_textbox(Inches(0.62), Inches(0.62), Inches(11), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Four searches, one shortlist"
    r.font.name, r.font.size, r.font.bold, r.font.color.rgb = SERIF, Pt(26), True, INK

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.20),
                                  Inches(12.09), Emu(9525))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RULE
    rule.line.fill.background()
    rule.shadow.inherit = False

    # ---- 1. the question --------------------------------------------------
    q = box(slide, 320, 8, 360, 44, fill=SEAL_WASH, stroke=SEAL)
    lines_in(q, [("A question", SANS, 13, SEAL, True)])

    line(slide, 500, 52, 500, 78, head=False)
    line(slide, 256, 78, 866, 78, head=False)
    for cx in (256, 622, 866):
        line(slide, cx, 78, cx, 98)

    # ---- 2. preparation ---------------------------------------------------
    prep = [
        (22, 468, 256, "Embed the question once",
         "one similarity search, read two different ways"),
        (510, 224, 622, "Widen the query  †",
         "synonyms + an invented passage"),
        (754, 224, 866, "Link the names",
         "then expand to every spelling"),
    ]
    for x, w, cx, title, sub in prep:
        b = box(slide, x, 100, w, 54, fill=CARD2, stroke=RULE)
        lines_in(b, [(title, SANS, 11, INK, True),
                     (sub, SANS, 8.5, INK_SOFT, False)])

    line(slide, 256, 154, 256, 172, head=False)
    line(slide, 134, 172, 378, 172, head=False)
    for cx in (134, 378):
        line(slide, cx, 172, cx, 196)
    line(slide, 622, 154, 622, 196)
    line(slide, 866, 154, 866, 196)

    # ---- 3. the four channels --------------------------------------------
    channels = [
        (22, 134, "vector", "dense cosine",
         "passages about the same idea,\nin different words"),
        (266, 378, "density", "document mass in the pool",
         "a document that keeps coming up\nis probably the right one"),
        (510, 622, "lexical", "BM25",
         "catches the exact term\nthe question used"),
        (754, 866, "graph", "entity → family → one hop",
         "follows the person, under\nevery spelling"),
    ]
    for x, cx, name, tech, plain in channels:
        b = box(slide, x, 198, 224, 104, fill=CARD2, stroke=RULE)
        rows = [(name, SANS, 12, INK, True), (tech, MONO, 8, SEAL, False)]
        rows += [(t, SANS, 8.5, INK_SOFT, False) for t in plain.split("\n")]
        lines_in(b, rows)
        line(slide, cx, 302, cx, 328)

    # ---- 4. the claim the whole diagram exists to make --------------------
    band = box(slide, 22, 332, 956, 36, fill=SEAL_WASH, stroke=SEAL_WASH, radius=0.02)
    lines_in(band, [("each ranks all 3,338 passages — nothing is narrowed before fusion",
                     SANS, 12, SEAL, True)])

    line(slide, 500, 368, 500, 398)

    # ---- 5. fusion --------------------------------------------------------
    f = box(slide, 270, 400, 460, 66, fill=SEAL_WASH, stroke=SEAL)
    lines_in(f, [("Rank fusion", SANS, 12.5, SEAL, True),
                 ("score = Σ 1 / (60 + rank)", MONO, 9, SEAL, False),
                 ("position only — no channel needs a weight", SANS, 8.5, INK_SOFT, False)])

    line(slide, 500, 466, 500, 498)

    # ---- 6. the eight slots ----------------------------------------------
    for i in range(4):
        b = box(slide, 106 + i * 100, 500, 88, 42, fill=OXIDE_WASH, stroke=OXIDE)
        lines_in(b, [("vector", SANS, 9.5, OXIDE, False)])
    for i in range(4):
        b = box(slide, 506 + i * 100, 500, 88, 42, fill=SEAL_WASH, stroke=SEAL)
        lines_in(b, [("fused", SANS, 9.5, SEAL, True)])

    label(slide, 300, 548, "four slots reserved, so fusion can only add",
          color=OXIDE, size=9)
    label(slide, 700, 548, "four slots decided by the vote", size=9)
    label(slide, 500, 572, "EIGHT PASSAGES", font=MONO, size=8, color=INK_FAINT)

    line(slide, 500, 594, 500, 624)

    # ---- 7. the answer ----------------------------------------------------
    a = box(slide, 270, 628, 460, 60, fill=CARD2, stroke=RULE)
    lines_in(a, [("Judge each passage on its own, then write", SANS, 12, INK, True),
                 ("one passage per call — never judged as a set",
                  SANS, 8.5, INK_SOFT, False),
                 ("if none of them helps, say so and stop",
                  SANS, 8.5, INK_SOFT, False)])

    # ---- speaker notes ----------------------------------------------------
    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "Why it is not a funnel. The usual design retrieves cheaply, then re-orders the "
        "survivors -- so anything ranked badly by the first stage is lost for good. Asked "
        "about obiter dicta, the passages that answer it sit at #329 and #2970 by "
        "similarity. Only a channel that reads the whole corpus on its own terms ever "
        "finds them."
    )
    for para in [
        "Vector and density share one embedding. They are not four independent opinions -- "
        "they are two independent searches plus two readings of the same similarity "
        "ranking. That shared origin is exactly why the reserved slots exist: channels "
        "correlated by construction cannot be fused with equal weight.",

        "† Widening the query is best-effort. Measured across eight demo questions, "
        "four came back unusable -- the model loops on synonyms until it runs out of room, "
        "always on the enumerative ones (“which courts appear...”). The lexical "
        "channel then searches the raw question instead, and now logs that it did.",

        "For the business half of the room: four opinions, combined by rank rather than by "
        "score, so nobody has to invent how much each one is worth. And when several "
        "channels independently pick the same passage, that agreement is the strongest "
        "signal the system produces.",
    ]:
        notes.add_paragraph().text = para

    prs.save(path)
    print("1 slide ->", path)
    print("  shapes:", len(slide.shapes), " notes:",
          len(slide.notes_slide.notes_text_frame.text), "chars +",
          len(slide.notes_slide.notes_text_frame.paragraphs) - 1, "paragraphs")


if __name__ == "__main__":
    build(sys.argv[1])
