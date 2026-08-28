#!/usr/bin/env python3
"""
Convert the interview deck (HTML) into an editable .pptx.

Written as a CONVERTER rather than a transcription so it can be re-run after
the HTML changes -- the HTML stayed the source of truth through many rounds of
edits, and hand-copying 16 slides into python-pptx would have forked it on the
first revision.

What it does NOT try to do: reproduce the web design. PowerPoint has no CSS
cascade, no auto-fit grid and none of the deck's three typefaces installed, so
chasing pixel fidelity would produce a brittle file that is unpleasant to edit.
The goal is a clean, on-palette, fully editable deck: real text boxes, real
tables, real speaker notes. Slimming it down afterwards is the point.

Usage:
    python/.venv/bin/python python/deck/html_to_pptx.py \
        <input.html> <output.pptx>
"""

import re
import sys
from html.parser import HTMLParser

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- palette
# Lifted from the HTML's :root block so the two stay recognisably the same
# deck. Light theme only -- a pptx has no viewer theme to respond to.
PAPER = RGBColor(0xF1, 0xF0, 0xEC)
CARD = RGBColor(0xFB, 0xFA, 0xF8)
INK = RGBColor(0x14, 0x17, 0x1B)
INK_SOFT = RGBColor(0x4A, 0x51, 0x58)
INK_FAINT = RGBColor(0x83, 0x8B, 0x93)
RULE = RGBColor(0xD8, 0xD5, 0xCD)
SEAL = RGBColor(0x1E, 0x3A, 0x5F)
SEAL_WASH = RGBColor(0xE7, 0xEB, 0xF1)
OXIDE = RGBColor(0xA6, 0x38, 0x2C)
BRASS = RGBColor(0x8A, 0x6D, 0x3B)

# Spectral / Public Sans / IBM Plex Mono are not installed on this machine and
# would silently substitute. Georgia + Helvetica Neue + Menlo are present on
# macOS and degrade predictably on Windows (Georgia ships there too; Helvetica
# Neue falls back to Arial).
SERIF = "Georgia"
SANS = "Helvetica Neue"
MONO = "Menlo"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.62)
CONTENT_W = SLIDE_W - 2 * MARGIN


# ================================================================ parsing
class Block:
    """One renderable thing on a slide."""

    def __init__(self, kind, **kw):
        self.kind = kind
        self.__dict__.update(kw)

    def __repr__(self):
        return f"<{self.kind}>"


class DeckParser(HTMLParser):
    """
    Pulls structure, not styling, out of the deck.

    Deliberately tolerant: it tracks only the handful of class names the deck
    actually uses to mean something (eyebrow, lede, note, fill, stage, stat,
    fix, notes) and treats everything else as inline text.
    """

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.slides = []
        self.cur = None
        self.stack = []          # (tag, classes)
        self.text = []           # inline text accumulator
        self.in_notes = False
        self.notes_buf = []
        self.list_items = []
        self.table = None
        self.row = None
        self.cell = None
        self.pending = None      # Block being assembled
        self.stats = []
        self.stat_v = None

    # -- helpers ---------------------------------------------------------
    def classes(self, attrs):
        d = dict(attrs)
        return set((d.get("class") or "").split())

    def flush_text(self):
        t = re.sub(r"\s+", " ", "".join(self.text)).strip()
        self.text = []
        return t

    def add(self, block):
        if self.cur is not None:
            self.cur["blocks"].append(block)

    def open_tags(self):
        return [t for t, _ in self.stack]

    def in_class(self, name):
        return any(name in c for _, c in self.stack)

    # -- handlers --------------------------------------------------------
    def handle_starttag(self, tag, attrs):
        cls = self.classes(attrs)
        self.stack.append((tag, cls))

        if tag == "br":
            self.text.append(" ")
            return

        if tag == "section" and "slide" in cls:
            self.cur = {"num": "", "eyebrow": "", "warn": False, "title": "",
                        "blocks": [], "notes": []}
            return
        if self.cur is None:
            return

        if tag == "details" and "notes" in cls:
            self.in_notes = True
            self.notes_buf = []
            self.text = []
            return

        if tag in ("ul", "ol") and not self.in_notes:
            self.list_items = []
        elif tag == "table" and not self.in_notes:
            self.table = {"head": [], "rows": []}
        elif tag == "tr" and self.table is not None:
            self.row = []
        elif tag in ("td", "th") and self.table is not None:
            self.cell = tag
            self.text = []
        elif tag == "div" and "stat" in cls:
            self.stats = []
            self.stat_v = None
            self.text = []
        elif tag == "span" and "iow" in cls:
            self.text.append(IOW)
        elif tag in ("p", "h1", "h2", "h3", "li", "div", "span", "summary"):
            if tag in ("p", "h1", "h2", "h3", "li", "summary") or cls & {
                "eyebrow", "num", "stage", "kv", "foot", "v", "k"
            }:
                self.text = []

    def handle_endtag(self, tag):
        if not self.stack:
            return
        _, cls = self.stack.pop()

        if tag == "section":
            if self.cur:
                self.slides.append(self.cur)
            self.cur = None
            return
        if self.cur is None:
            return

        if tag == "details" and self.in_notes:
            self.in_notes = False
            self.cur["notes"] = [p for p in self.notes_buf if p]
            return

        if self.in_notes:
            if tag == "p":
                t = self.flush_text()
                if t and t.lower() != "presenter notes":
                    self.notes_buf.append(t)
            elif tag == "summary":
                self.flush_text()
            return

        # ---- tables
        if tag in ("td", "th") and self.table is not None:
            (self.row if self.row is not None else []).append(self.flush_text())
            self.cell = None
            return
        if tag == "tr" and self.table is not None:
            if self.row:
                if all(c == "" for c in self.row):
                    pass
                elif not self.table["head"] and "thead" in self.open_tags() + ["thead"] \
                        and self.table["rows"] == [] and self._last_was_header:
                    self.table["head"] = self.row
                else:
                    self.table["rows"].append(self.row)
            self.row = None
            return
        if tag == "table" and self.table is not None:
            self.add(Block("table", head=self.table["head"], rows=self.table["rows"]))
            self.table = None
            return

        # ---- lists
        if tag == "li":
            t = self.flush_text()
            if t:
                self.list_items.append(t)
            return
        if tag in ("ul", "ol"):
            if self.list_items:
                self.add(Block("bullets", items=self.list_items))
            self.list_items = []
            return

        # ---- headings and text
        if tag == "h1":
            self.cur["title"] = self.flush_text()
            return
        if tag == "h2":
            self.cur["title"] = self.flush_text()
            return
        if tag == "h3":
            t = self.flush_text()
            if t:
                self.add(Block("h3", text=t, warn="color:var(--oxide)" in str(cls)))
            return
        if tag == "p":
            t = self.flush_text()
            if not t:
                return
            if "lede" in cls:
                self.add(Block("lede", text=t))
            elif "fix" in cls:
                self.add(Block("fix", text=t))
            elif "sub" in cls:
                self.add(Block("sub", text=t))
            else:
                self.add(Block("body", text=t, note=self.in_class("note"),
                               warn=self.in_class("warn")))
            return
        if tag == "div":
            t = self.flush_text()
            if "eyebrow" in cls:
                self.cur["eyebrow"] = t
                self.cur["warn"] = "warn" in cls
            elif "num" in cls:
                self.cur["num"] = t
            elif "stage" in cls and t:
                self.add(Block("stage", text=t))
            elif "foot" in cls and t:
                self.add(Block("foot", text=t))
            elif "v" in cls:
                # The big number. Held until its label arrives -- they are two
                # sibling divs, and a number with no label says nothing.
                self.stat_v = (t, "bad" in cls)
            elif "k" in cls and self.stat_v is not None:
                self.stats.append((self.stat_v[0], t, self.stat_v[1]))
                self.stat_v = None
            elif "stat" in cls and self.stats:
                self.add(Block("stats", pairs=self.stats))
                self.stats = []
            return
        if tag == "span":
            if "num" in cls and self.cur is not None and not self.cur["num"]:
                self.cur["num"] = self.flush_text()
            elif "n" in cls:
                # The step number in a .stage row is a separate span with no
                # whitespace between it and the label, so "01" and "Split" ran
                # together once the tags were stripped.
                self.text.append("  ")
            return

    def handle_startendtag(self, tag, attrs):
        # <br /> only; the bare <br> form arrives via handle_starttag, so both
        # are handled -- the deck's h1 uses one to break the title line and
        # without this the two halves ran together.
        if tag == "br":
            self.text.append(" ")

    def handle_data(self, data):
        self.text.append(data)

    @property
    def _last_was_header(self):
        return self.cell == "th" or True


def parse_deck(path):
    p = DeckParser()
    p.feed(open(path, encoding="utf-8").read())
    return p.slides


# ================================================================ rendering
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# Separates a bullet's technical statement from its plain-language restatement.
# The deck writes the second half as <span class="iow">In other words: ...</span>;
# the parser drops this marker in so the renderer can split them back apart and
# set the restatement on its own indented italic line instead of running the two
# sentences together in one bullet.
IOW = "\u241f"


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return box, tf


def style(run, *, font=SANS, size=14, color=INK, bold=False, italic=False,
          spacing=None):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def para(tf, text, *, first=False, font=SANS, size=14, color=INK, bold=False,
         italic=False, space_before=6, space_after=4, bullet=False,
         indent=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if indent:
        p.level = indent
    r = p.add_run()
    r.text = ("•  " if bullet else "") + text
    style(r, font=font, size=size, color=color, bold=bold, italic=italic)
    return p


def add_rule(slide, top):
    """The thin horizontal line under the eyebrow -- the deck's one repeated
    graphic device, and cheap to reproduce."""
    from pptx.enum.shapes import MSO_SHAPE
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, top,
                                CONTENT_W, Emu(9525))  # ~0.75pt
    ln.fill.solid()
    ln.fill.fore_color.rgb = RULE
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def render_table(slide, block, left, top, width):
    """
    A real pptx table, with row heights estimated from wrapped content.

    The estimate matters: PowerPoint grows a row to fit its text regardless of
    the height we ask for, so a naive fixed height per row makes this function
    under-report where the table ends -- and the next block lands on top of it.
    That is exactly what happened on the positioning slide, where two cells
    wrap to three lines.
    """
    rows = len(block.rows) + (1 if block.head else 0)
    cols = max([len(r) for r in block.rows] + [len(block.head or [])])
    if rows == 0 or cols == 0:
        return top

    col_w = int(width / cols)
    col_w_in = col_w / Inches(1)

    def lines_for(texts, size):
        chars_per_line = max(12, int(col_w_in * 72 / (size * 0.52)))
        return max(1, max(-(-len(t) // chars_per_line) for t in texts))

    heights = []
    if block.head:
        heights.append(Inches(0.13 + 0.16 * lines_for(block.head, 10)))
    for r in block.rows:
        heights.append(Inches(0.14 + 0.185 * lines_for(r or [""], 11)))
    total = sum(heights, Emu(0))

    shape = slide.shapes.add_table(rows, cols, left, top, width, total)
    tbl = shape.table
    tbl.first_row = bool(block.head)
    for c in range(cols):
        tbl.columns[c].width = col_w
    for i, h in enumerate(heights):
        tbl.rows[i].height = h

    def fill(cell, text, header):
        cell.margin_left = cell.margin_right = Inches(0.07)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.TOP
        cell.fill.solid()
        cell.fill.fore_color.rgb = SEAL_WASH if header else CARD
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        style(r, size=10 if header else 11,
              color=INK_SOFT if header else INK, bold=header)

    r0 = 0
    if block.head:
        for c, txt in enumerate(block.head):
            fill(tbl.cell(0, c), txt, True)
        r0 = 1
    for i, row in enumerate(block.rows):
        for c in range(cols):
            fill(tbl.cell(i + r0, c), row[c] if c < len(row) else "", False)
    return top + total + Inches(0.16)


def est_height(blocks, width_in):
    """
    Rough height of a run of text blocks, in inches.

    PowerPoint has no layout engine we can query from here, so this estimates
    wrapped line count from character count. It only has to be good enough to
    stop a table landing on top of a paragraph -- everything is a real text box
    and can be nudged by hand afterwards.
    """
    SIZES = {"lede": 14, "h3": 11, "bullets": 12, "stage": 12, "fix": 11,
             "sub": 9.5, "foot": 9, "stats": 20, "body": 12}
    total = 0.0
    for b in blocks:
        size = SIZES.get(b.kind, 12)
        if b.kind == "bullets":
            texts = b.items
        elif b.kind == "stats":
            texts = [f"{v} {k}" for v, k, _ in b.pairs]
        else:
            texts = [getattr(b, "text", "")]
        for t in texts:
            chars_per_line = max(20, int(width_in * 72 / (size * 0.5)))
            lines = max(1, -(-len(t) // chars_per_line))
            total += lines * (size * 1.3) / 72 + 0.07
            # A bullet carrying a restatement renders as two paragraphs, so it
            # costs one more inter-paragraph gap than its length implies.
            if IOW in t:
                total += 0.09
    return total


def render_text_group(slide, blocks, top, width):
    """One textbox holding a consecutive run of non-table blocks."""
    height = Inches(est_height(blocks, width / Inches(1)) + 0.1)
    box, tf = textbox(slide, MARGIN, top, width, height)
    first = True
    for b in blocks:
        if b.kind == "lede":
            para(tf, b.text, first=first, size=14, color=INK_SOFT,
                 space_before=0 if first else 8)
        elif b.kind == "h3":
            para(tf, b.text.upper(), first=first, font=SANS, size=11,
                 color=OXIDE if getattr(b, "warn", False) else INK_SOFT,
                 bold=True, space_before=0 if first else 12, space_after=2)
        elif b.kind == "bullets":
            for i, it in enumerate(b.items):
                main, _, iow = it.partition(IOW)
                para(tf, main.strip(), first=first and i == 0, size=12, color=INK,
                     bullet=True, space_before=0 if (first and i == 0) else 3,
                     space_after=1 if iow.strip() else 3)
                if iow.strip():
                    para(tf, iow.strip(), font=SERIF, size=11, color=INK_SOFT,
                         italic=True, indent=1, space_before=0, space_after=4)
        elif b.kind == "stage":
            para(tf, b.text, first=first, size=12, color=INK,
                 space_before=0 if first else 5, space_after=3)
        elif b.kind == "fix":
            para(tf, "\u2192  " + b.text, first=first, size=11, color=SEAL,
                 space_before=2, space_after=6)
        elif b.kind == "stats":
            for i, (val, key, bad) in enumerate(b.pairs):
                pp = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
                pp.space_before = Pt(0 if (first and i == 0) else 6)
                pp.space_after = Pt(2)
                rv = pp.add_run()
                rv.text = val + "   "
                style(rv, font=SERIF, size=20,
                      color=OXIDE if bad else SEAL, bold=True)
                rk = pp.add_run()
                rk.text = key.upper()
                style(rk, font=SANS, size=10, color=INK_FAINT)
        elif b.kind == "sub":
            para(tf, b.text, first=first, size=9.5, color=INK_FAINT,
                 space_before=6)
        elif b.kind == "foot":
            para(tf, b.text, first=first, font=MONO, size=9, color=INK_FAINT,
                 space_before=10)
        else:  # body
            warn = getattr(b, "warn", False)
            note = getattr(b, "note", False)
            para(tf, b.text, first=first, size=12,
                 color=OXIDE if warn else (SEAL if note else INK),
                 italic=note and not warn,
                 space_before=0 if first else 7)
        first = False
    return top + height


def render_slide(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, CARD)

    y = MARGIN

    # slide number, top right
    box, tf = textbox(slide, SLIDE_W - MARGIN - Inches(1.0), Inches(0.3),
                      Inches(1.0), Inches(0.3))
    p = para(tf, s["num"], first=True, font=MONO, size=9, color=INK_FAINT,
             space_before=0, space_after=0)
    p.alignment = PP_ALIGN.RIGHT

    if s["eyebrow"]:
        box, tf = textbox(slide, MARGIN, y, CONTENT_W, Inches(0.25))
        para(tf, s["eyebrow"].upper(), first=True, font=MONO, size=9,
             color=OXIDE if s["warn"] else SEAL, space_before=0, space_after=0)
        y += Inches(0.3)

    if s["title"]:
        title_size = 34 if s["num"] == "00" else 26
        # Height has to follow the wrap, not a constant. A title long enough to
        # take two lines used to overrun the rule drawn underneath it, and the
        # rule was painted straight through the second line.
        if s["num"] == "00":
            h = Inches(0.9)
        else:
            cpl = max(20, int((CONTENT_W / Inches(1)) * 72 / (title_size * 0.56)))
            lines = max(1, -(-len(s["title"]) // cpl))
            h = Inches(0.55 * lines)
        box, tf = textbox(slide, MARGIN, y, CONTENT_W, h)
        para(tf, s["title"], first=True, font=SERIF, size=title_size,
             color=INK, bold=True, space_before=0, space_after=0)
        y += h + Inches(0.08)

    add_rule(slide, y)
    y += Inches(0.2)

    # Blocks flow top to bottom in document order. Consecutive text blocks
    # share one textbox so their paragraph spacing behaves; a table breaks the
    # run, gets its own shape, and the cursor resumes below it. The previous
    # version placed all tables at an estimated offset, which put the note on
    # slide 03 underneath its own table.
    group = []
    for b in s["blocks"]:
        if b.kind == "table":
            if group:
                y = render_text_group(slide, group, y, CONTENT_W)
                group = []
            y = render_table(slide, b, MARGIN, y, CONTENT_W)
        else:
            group.append(b)
    if group:
        render_text_group(slide, group, y, CONTENT_W)

    if s["notes"]:
        tfn = slide.notes_slide.notes_text_frame
        tfn.text = s["notes"][0]
        for n in s["notes"][1:]:
            tfn.add_paragraph().text = n

    return slide


def main(src, dst):
    slides = parse_deck(src)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for s in slides:
        render_slide(prs, s)
    prs.save(dst)
    print(f"{len(slides)} slides -> {dst}")
    for s in slides:
        print(f"  {s['num']:>3}  {s['title'][:54]:<54} "
              f"{len(s['blocks']):>2} blocks, {len(s['notes'])} notes")


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
