"""Generates the technical-walkthrough .pptx for the Graph RAG / ArangoDB project.

Content is drawn from _project_step_by_step_plan.md (the pipeline design and
build log) and _step3_extraction_reference.md (the extraction prompt detail).
Re-run this script any time to regenerate the deck from scratch -- it does not
read an existing .pptx, so manual edits made directly in PowerPoint/Slides
will be lost on the next run. Layout helpers live in deck_theme.py, shared
with build_deck_story.py so both decks stay visually consistent.
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from deck_theme import (
    NAVY, TEAL, LIGHT_TEAL, DARK_TEXT, WHITE, GRAY, LIGHT_BLUE_TEXT,
    new_presentation, add_slide, set_background, add_textbox, add_bullets,
    add_header, add_footer, add_box, add_plain_box, add_arrow, add_callout,
)

prs, BLANK = new_presentation()

# ============================================================================
# Slide 1 -- Title
# ============================================================================
s = add_slide(prs, BLANK)
set_background(s, NAVY)
add_textbox(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.2),
            "Graph RAG on ArangoDB", size=44, color=WHITE, bold=True)
add_textbox(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.8),
            "Extracting a Knowledge Graph from 187 Real Court Filings",
            size=22, color=TEAL, bold=False)
add_textbox(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.6),
            "SDNY Docket 447706 — A Technical Walkthrough", size=16, color=LIGHT_BLUE_TEXT, italic=True)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.45), Inches(2.2), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = TEAL
line.line.fill.background()
line.shadow.inherit = False

# ============================================================================
# Slide 2 -- The Problem
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "The Problem", kicker="Why this project exists")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5),
            [
                "187 real, unsealed PDF court filings (~182 MB), SDNY docket 447706 — Epstein-litigation-related",
                "Confirmed as real text-layer PDFs (no OCR needed) via a pypdf extraction test",
                "Content is unstructured prose: names, roles, relationships, dates, and docket numbers are buried in filing text, not tagged data",
                "Goal: turn 187 documents into a queryable knowledge graph plus a corpus an LLM can answer questions against, grounded in citations",
                "Explicit constraint from day one: local-vs-cloud LLM host is a deliberate decision, not a default — the corpus likely names non-party individuals (victims/witnesses)",
            ], size=18, space_after=18)
add_footer(s, 2)

# ============================================================================
# Slide 3 -- Architecture at a Glance
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Architecture at a Glance", kicker="Two databases, two languages, one local model")

col_w = Inches(5.6)
col_h = Inches(3.0)
top = Inches(1.7)

add_plain_box(s, Inches(0.7), top, col_w, col_h, NAVY)
add_plain_box(s, Inches(6.9), top, col_w, col_h, TEAL)

add_textbox(s, Inches(1.0), Inches(1.85), col_w - Inches(0.6), Inches(0.4),
            "ArangoDB", size=20, color=WHITE, bold=True)
add_textbox(s, Inches(1.0), Inches(2.25), col_w - Inches(0.6), Inches(0.3),
            "Graph store", size=13, color=LIGHT_BLUE_TEXT, italic=True)
add_bullets(s, Inches(1.0), Inches(2.65), col_w - Inches(0.6), Inches(2),
            [
                "je_chunks — chunk text + metadata",
                "je_entities, je_relations (planned)",
                "je_mentioned_in edges (planned)",
            ], size=14, color=WHITE, space_after=8)

add_textbox(s, Inches(7.2), Inches(1.85), col_w - Inches(0.6), Inches(0.4),
            "LanceDB", size=20, color=WHITE, bold=True)
add_textbox(s, Inches(7.2), Inches(2.25), col_w - Inches(0.6), Inches(0.3),
            "Vector store", size=13, color=LIGHT_TEAL, italic=True)
add_bullets(s, Inches(7.2), Inches(2.65), col_w - Inches(0.6), Inches(2),
            [
                "je_chunks — chunk-level embeddings",
                "je_entity_mentions (planned) — per-mention context-snippet embeddings",
            ], size=14, color=WHITE, space_after=8)

add_textbox(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(0.4),
            "Language split", size=16, color=NAVY, bold=True)
add_bullets(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.8), [
    "Node.js — Express web layer, upload UI, SSE progress streaming, MCP tool delivery",
    "Python — PDF text extraction and chunking (pypdf)",
    "Local LLM — qwen3:8b via Ollama, used for extraction and dedup judgment calls",
], size=15)
add_footer(s, 3)

# ============================================================================
# Slide 4 -- The Pipeline, End to End (diagram)
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "The Pipeline, End to End", kicker="Steps 1–8, repeated per chunk / per file")

steps = [
    ("1", "Intake", "chunk · embed · store"),
    ("3", "Extract", "regex + LLM"),
    ("4", "Candidates", "Jaro-Winkler + embed"),
    ("5", "Judge", "LLM: same entity?"),
    ("6", "Merge", "write graph"),
    ("8", "Query", "RAG at ask-time"),
]
box_w = Inches(1.8)
box_h = Inches(1.15)
arrow_w = Inches(0.3)
start_x = Inches(0.5)
y = Inches(2.4)
x = start_x
for i, (num, label, sub) in enumerate(steps):
    fill = NAVY if i % 2 == 0 else TEAL
    add_box(s, x, y, box_w, box_h, f"Step {num}: {label}", sub, fill=fill, size=13)
    x = Emu(x + box_w)
    if i < len(steps) - 1:
        add_arrow(s, x, Emu(y + box_h // 2 - Emu(Inches(0.12))), arrow_w, Inches(0.24))
        x = Emu(x + arrow_w)

add_bullets(s, Inches(0.7), Inches(4.1), Inches(11.9), Inches(2.6), [
    "Step 2 (acknowledgement) and Step 7 (repeat across all chunks/files) are operational, not pictured above",
    "Steps 1 and 3 are built and running against the real corpus today; Steps 4–6 are fully designed, not yet coded",
    "Step 8 (query-time RAG) is a later phase — not designed in detail yet",
    "Chunk-and-merge is the scaling strategy throughout: no step ever holds the whole corpus in one LLM context — external databases (ArangoDB / LanceDB) are the memory, not the prompt",
], size=15, space_after=10)
add_footer(s, 4)

# ============================================================================
# Slide 5 -- Step 1: Intake Pipeline
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Step 1 — Intake Pipeline", kicker="Deterministic, no LLM")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "Extract raw text from each PDF (Python, pypdf) — no OCR needed, real text-layer PDFs",
    "Split into overlapping chunks: 350 words per chunk, 50-word overlap — sized for the embedding model's 512-token ceiling",
    "Embed each chunk (bge-large via Ollama's /api/embed)",
    "Write chunk text + vector + metadata (source filename, chunk index) into both LanceDB (vector) and ArangoDB (graph-side record) — cross-verified to match exactly",
    "Hardened idempotency: re-uploading a file deletes that file's existing rows before re-adding (delete-then-add), so a repeat upload can never create duplicates",
    ("Status: complete and verified — 187 files → 3,575 chunks, confirmed matching across both databases", 0),
], size=17, space_after=14)
add_footer(s, 5)

# ============================================================================
# Slide 6 -- Step 3: Entity Extraction
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Step 3 — Entity Extraction", kicker="Deterministic pre-pass + one narrow LLM call")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "Pre-pass (regex, no LLM): docket/case numbers and dates follow fixed patterns — pulled out with zero hallucination risk",
    "LLM call (qwen3:8b, think disabled, JSON-schema-constrained output): judges only person / organization / court entities — the genuinely ambiguous part",
    "A fourth type, \"reference\", is an explicit discard bucket for citations to other cases/documents/exhibits — backed by an independent regex safety net checking the model's own output",
    "candidate_role is only assigned when the role word sits directly adjacent to the name in the text — never inferred from context or prior knowledge",
    "Relation extraction was deliberately deferred: asking for entities + relations in one call roughly quadrupled output length and caused fabricated relations",
], size=17, space_after=14)
add_footer(s, 6)

# ============================================================================
# Slide 7 -- Design Decision: Killing Hallucinations by Architecture
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Design Decision: Architecture Over Prompting", kicker="The “Judge LAP” hallucination")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(3.2), [
    "Early testing: the LLM read a docket number's judge-initials suffix (“...-LAP”) and invented a fictional person, “Judge LAP”",
    "The failure wasn't a mismatch — it was misinterpretation. No amount of prompt tuning fully closes that gap for a model that's allowed to interpret",
    "Fix: moved docket numbers (and dates) out of the LLM's job entirely, into deterministic regex extraction",
], size=17, space_after=16)
add_callout(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.6),
            "A regex can only match literal text — it has no capacity to “interpret” a suffix as a name. "
            "This failure mode is structurally impossible now, not just less likely.")
add_footer(s, 7)

# ============================================================================
# Slide 8 -- Design Decision: Accepting Non-Determinism
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Design Decision: Accepting Non-Determinism", kicker="Working on percentages, not guarantees")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "LLM sampling isn't fully deterministic even at temperature 0 (floating-point non-associativity in parallelized inference) — observed directly: the same chunk found 5 entities in one run, 3 in another",
    "Rejected: running extraction twice per chunk and taking the union — a real lever, but not adopted without evidence it's actually needed; paying 2x runtime on one anecdotal data point isn't justified yet",
    "Adopted instead: keep the LLM's job as narrow as possible (deterministic pre-pass shrinks what it has to judge), and let Step 5's “unsure → flag for human review, never auto-merge” absorb the remaining variance",
    "A false merge (conflating two real people) is judged worse than a missed merge — for legal-accuracy reasons",
    "This is a deliberate quality bar: this pipeline is not safety-critical, and probabilistic/percentage-based accuracy is an accepted, explicit tradeoff — not a gap to be engineered away",
], size=17, space_after=14)
add_footer(s, 8)

# ============================================================================
# Slide 9 -- Steps 4-6: Dedup and Merge
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Steps 4–6 — Dedup Candidates, Judgment, Merge", kicker="Designed in full; not yet built")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "Step 4 (deterministic): two independent candidate channels, unioned — Jaro-Winkler string similarity, plus embedding similarity over per-mention context snippets (not bare names, which don't carry enough signal alone)",
    "je_entity_mentions (LanceDB, planned): one row per mention, not per entity — preserves phrasing diversity across 187 documents instead of blurring it into one averaged vector",
    "Step 5 (LLM judgment): shown current context + each candidate's known facts, decides same / different / genuinely unsure — unsure never auto-merges",
    "Step 6 (deterministic write): confirmed match upserts onto the existing je_entities node; new entity gets a new node; ambiguous is written with a review flag",
    "Write-order rule: entity_id must exist before anything references it — entity create/upsert always happens before je_mentioned_in edges or je_entity_mentions vector rows are written",
], size=16, space_after=13)
add_footer(s, 9)

# ============================================================================
# Slide 10 -- Step 8: Query-Time RAG
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Step 8 — Query-Time RAG", kicker="Later phase, not designed in detail yet")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "User asks a plain-English question",
    "Vector search retrieves relevant chunks (LanceDB) + graph traversal pulls related entities (ArangoDB)",
    "LLM composes a cited answer grounded in both — never the whole corpus in context, only what's relevant to the question",
    "Same chunk-and-merge principle as ingestion: the corpus's “memory” lives in the databases, retrieval happens at ask-time, not by loading everything upfront",
    "Model choice for retrieval/answer-composition is intentionally separate from the extraction model — parked for a future side-by-side comparison",
], size=18, space_after=16)
add_footer(s, 10)

# ============================================================================
# Slide 11 -- Current Status
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "Current Status", kicker="As of this deck")
stats = [
    ("187", "PDF files ingested"),
    ("3,575", "chunks, verified matching\nacross ArangoDB + LanceDB"),
    ("714 / 3,575", "chunks entity-extracted\n(batch run in progress)"),
]
box_w = Inches(3.7)
gap = Inches(0.4)
x = Inches(0.7)
for val, lbl in stats:
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), box_w, Inches(1.9))
    shp.fill.solid()
    shp.fill.fore_color.rgb = NAVY
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = val
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = TEAL
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = lbl
    r2.font.size = Pt(13)
    r2.font.color.rgb = WHITE
    x = Emu(x + box_w + gap)

add_bullets(s, Inches(0.7), Inches(4.1), Inches(11.9), Inches(2.6), [
    "Step 1 (intake): complete and verified against the full real corpus",
    "Step 3 (extraction): fully built and proven; the full 3,575-chunk batch run is in progress",
    "Steps 4–6 (dedup, judgment, merge): fully designed, zero code written yet",
    "Step 8 (query-time RAG): later phase, not designed in detail yet",
], size=16, space_after=10)
add_footer(s, 11)

# ============================================================================
# Slide 12 -- What's Next
# ============================================================================
s = add_slide(prs, BLANK)
add_header(s, "What's Next", kicker="Closing")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5), [
    "Finish the full-corpus extraction batch run, verify results",
    "Build Step 4: candidate generation (Jaro-Winkler + je_entity_mentions embeddings)",
    "Build Step 5: LLM dedup judgment call",
    "Build Step 6: merge/write to je_entities, je_relations, je_mentioned_in",
    "Phase 2 (parked): human-review web page for flagged possible-duplicates — async, not a live chat interruption during batch runs",
    "Later: relation extraction as its own pass; model comparison for the query-time RAG role",
], size=18, space_after=16)
add_footer(s, 12)

# ============================================================================
out_path = "/Users/budubasa/Documents/N_Claude_Projects/arango_graph_rag_llm/Graph_RAG_LLM_Overview.pptx"
prs.save(out_path)
print(f"Saved {out_path} ({len(prs.slides)} slides)")
