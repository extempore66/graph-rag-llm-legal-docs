"""Shared visual house-style for this project's .pptx decks (python-pptx helpers).

Used by build_deck.py (technical walkthrough) and build_deck_story.py (article
companion). Keeping the palette/layout helpers here means both decks stay
visually consistent and a style tweak only has to be made once.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Palette ---------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x2E, 0x86, 0xAB)
LIGHT_TEAL = RGBColor(0xE8, 0xF3, 0xF7)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_BLUE_TEXT = RGBColor(0xC8, 0xD4, 0xE8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    """Returns (prs, blank_layout) -- call once per deck script."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs, prs.slide_layouts[6]


def add_slide(prs, blank_layout):
    return prs.slides.add_slide(blank_layout)


def set_background(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, color=DARK_TEXT,
                 bold=False, align=PP_ALIGN.LEFT, font="Calibri", italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=DARK_TEXT,
                 space_after=10, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(space_after)
        run = p.add_run()
        bullet_char = "•  " if level == 0 else "–  "
        run.text = bullet_char + text
        run.font.size = Pt(size if level == 0 else size - 2)
        run.font.color.rgb = color if level == 0 else GRAY
        run.font.name = font
    return box


def add_header(slide, title, kicker=None):
    set_background(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False
    if kicker:
        add_textbox(slide, Inches(0.5), Inches(0.12), Inches(11), Inches(0.3),
                    kicker.upper(), size=11, color=TEAL, bold=True)
        add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.65),
                    title, size=28, color=WHITE, bold=True)
    else:
        add_textbox(slide, Inches(0.5), Inches(0.28), Inches(12), Inches(0.65),
                    title, size=28, color=WHITE, bold=True)
    return bar


def add_footer(slide, page_num, label="Graph RAG on ArangoDB"):
    add_textbox(slide, Inches(0.5), Inches(7.15), Inches(6), Inches(0.3),
                label, size=10, color=GRAY)
    add_textbox(slide, Inches(12.3), Inches(7.15), Inches(0.6), Inches(0.3),
                str(page_num), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_plain_box(slide, left, top, width, height, fill, line_color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    shp.text_frame.paragraphs[0].text = ""
    return shp


def add_box(slide, left, top, width, height, label, sublabel=None,
            fill=TEAL, text_color=WHITE, size=13):
    """Centered label + optional sublabel inside a rounded box (no separate bullets)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = WHITE
    shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    if sublabel:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sublabel
        r2.font.size = Pt(size - 3)
        r2.font.color.rgb = text_color
    return shp


def add_arrow(slide, left, top, width, height, horizontal=True):
    shape_type = MSO_SHAPE.RIGHT_ARROW if horizontal else MSO_SHAPE.DOWN_ARROW
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = GRAY
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_callout(slide, left, top, width, height, text, size=18):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_TEAL
    box.line.color.rgb = TEAL
    box.line.width = Pt(1.25)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(18)
    tf.margin_right = Pt(18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.italic = True
    run.font.color.rgb = NAVY
    return box
